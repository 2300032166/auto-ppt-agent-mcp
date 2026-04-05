# mcp_servers/ppt_server.py
# ─────────────────────────────────────────────────────────────────────────────
# PPT MCP Server  —  with full THEME CUSTOMIZATION support
#
# Tools exposed:
#   set_theme(theme_name)                       NEW — must call BEFORE create
#   create_presentation(title)
#   add_slide(title)
#   write_text(slide_index, points)
#   save_presentation(filename)
#
# Theme pipeline (server-side):
#   1. set_theme()        → stores ThemeConfig in module state
#   2. create_presentation() → applies theme background + styles title slide
#   3. add_slide()        → applies theme background + styles content slide
#   4. write_text()       → applies theme font/colour to bullet text
# ─────────────────────────────────────────────────────────────────────────────

import os
import sys
import re
import logging
from typing import List, Annotated, Optional

from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pydantic import BaseModel, Field

# ── Bootstrap path so themes/ is importable when server runs standalone ────
_BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if _BASE_DIR not in sys.path:
    sys.path.insert(0, _BASE_DIR)

from themes.theme_config import ThemeConfig, get_theme, DEFAULT_THEME_NAME

# ── Bullet label-prefix cleaner ────────────────────────────────────────────
# Last-defence guard: strips any [Tag] or (Tag) prefix that may arrive via
# any code path, so nothing unprofessional ever lands on a slide.
_LABEL_RE = re.compile(r"^[\[(][^\]\)]{1,50}[\])]\s*", re.IGNORECASE)

# ── Logging ────────────────────────────────────────────────────────────────
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("PPT-Server")

# ── FastMCP instance ───────────────────────────────────────────────────────
mcp = FastMCP("PPT-Server")

# ── Module-level state ─────────────────────────────────────────────────────
_current_presentation: Optional[Presentation] = None
_active_theme: ThemeConfig = get_theme(DEFAULT_THEME_NAME)   # always non-None
_output_path: str = os.path.join(_BASE_DIR, "output")


# ═══════════════════════════════════════════════════════════════════════════
# Pydantic response models
# ═══════════════════════════════════════════════════════════════════════════

class ToolResult(BaseModel):
    result: str = Field(..., description="A success or error message")

class SlideToolResult(ToolResult):
    slide_index: int = Field(..., description="The index of the slide affected")


# ═══════════════════════════════════════════════════════════════════════════
# Internal helpers  (NOT MCP tools)
# ═══════════════════════════════════════════════════════════════════════════

def _rgb(color_tuple) -> RGBColor:
    """Convert an (R, G, B) tuple to an RGBColor."""
    r, g, b = color_tuple
    return RGBColor(r, g, b)


def _apply_background(slide, theme: ThemeConfig) -> None:
    """
    Set a solid background colour on *slide* using a filled background shape.

    python-pptx's Background API is fragile across template types, so we use
    a full-slide rectangle at z-order 0 as a reliable fallback.
    """
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    from lxml import etree

    # Use the native slide background fill when possible
    try:
        from pptx.oxml.ns import nsmap
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(theme.background_color)
        return
    except Exception:
        pass

    # Fallback: insert a coloured rectangle at the very back
    slide_width  = slide.shapes._spTree.getparent().getparent().cSld.attrib.get("cx", 9144000)
    slide_height = slide.shapes._spTree.getparent().getparent().cSld.attrib.get("cy", 5143500)
    try:
        slide_width  = int(slide_width)
        slide_height = int(slide_height)
    except (ValueError, TypeError):
        slide_width, slide_height = 9144000, 5143500

    rect = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, slide_width, slide_height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = _rgb(theme.background_color)
    rect.line.fill.background()  # no border

    # Move the rectangle to the back (first child of spTree)
    sp_tree = slide.shapes._spTree
    sp_el   = rect._element
    sp_tree.remove(sp_el)
    sp_tree.insert(2, sp_el)          # index 2 = after nvGrpSpPr and grpSpPr


def _apply_accent_bar(slide, theme: ThemeConfig) -> None:
    """
    Draw a thin coloured bar at the top of a content slide (under the title).
    Height: 6 px equivalent (~68 960 EMU).
    """
    try:
        prs_width = slide.shapes._spTree.getparent().getparent().slide_width
    except Exception:
        prs_width = Inches(10)

    bar_height = Pt(4)
    bar_top    = Inches(1.35)

    bar = slide.shapes.add_shape(1, 0, bar_top, prs_width, bar_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(theme.accent_color)
    bar.line.fill.background()


def _style_text_frame(shape, color: tuple, font_family: str, font_size: int,
                       bold: bool = False) -> None:
    """Apply colour / font / size to every paragraph run in *shape*."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = _rgb(color)
            run.font.name      = font_family
            run.font.size      = Pt(font_size)
            run.font.bold      = bold
        # Also set default on the paragraph-level font
        para.font.color.rgb = _rgb(color)
        para.font.name      = font_family
        para.font.size      = Pt(font_size)
        para.font.bold      = bold


def _theme_title_slide(slide, theme: ThemeConfig) -> None:
    """
    Apply theme to the opening (title) slide:
      • background
      • title text → title_color / title_size
      • subtitle   → subtitle_color / body_size
    """
    _apply_background(slide, theme)

    for shape in slide.placeholders:
        ph_type = shape.placeholder_format.type
        if ph_type == PP_PLACEHOLDER.TITLE or ph_type == PP_PLACEHOLDER.CENTER_TITLE:
            _style_text_frame(shape, theme.title_color,
                              theme.font_family, theme.title_size, bold=True)
        elif ph_type == PP_PLACEHOLDER.SUBTITLE:
            _style_text_frame(shape, theme.subtitle_color,
                              theme.font_family, theme.body_size)


def _theme_content_slide(slide, theme: ThemeConfig) -> None:
    """
    Apply theme to a content slide:
      • background
      • accent bar
      • slide title → slide_title_color / slide_title_size
      • body text   → body_color / body_size
    """
    _apply_background(slide, theme)
    _apply_accent_bar(slide, theme)

    for shape in slide.placeholders:
        ph_type = shape.placeholder_format.type
        if ph_type == PP_PLACEHOLDER.TITLE:
            _style_text_frame(shape, theme.slide_title_color,
                              theme.font_family, theme.slide_title_size, bold=True)
        else:
            _style_text_frame(shape, theme.body_color,
                              theme.font_family, theme.body_size)


# ═══════════════════════════════════════════════════════════════════════════
# MCP TOOL: set_theme
# ═══════════════════════════════════════════════════════════════════════════

@mcp.tool()
def set_theme(theme_name: str = "professional") -> ToolResult:
    """
    Set the global presentation theme.

    Must be called BEFORE create_presentation().

    Args:
        theme_name: One of 'professional', 'dark', 'academic', 'creative',
                    'minimal'.  Unknown values fall back to 'professional'.

    Returns:
        ToolResult confirming the active theme.
    """
    global _active_theme
    _active_theme = get_theme(theme_name)
    msg = (
        f"Theme set to '{_active_theme.name}' | "
        f"bg={_active_theme.background_color} | "
        f"font={_active_theme.font_family} | "
        f"title_size={_active_theme.title_size}pt | "
        f"body_size={_active_theme.body_size}pt"
    )
    logger.info(f"[THEME] {msg}")
    return ToolResult(result=msg)


# ═══════════════════════════════════════════════════════════════════════════
# MCP TOOL: create_presentation
# ═══════════════════════════════════════════════════════════════════════════

@mcp.tool()
def create_presentation(title: str) -> ToolResult:
    """
    Creates a new PowerPoint presentation with a themed title slide.

    Inherits the theme set by set_theme(); defaults to 'professional' if
    set_theme() was never called.
    """
    global _current_presentation, _active_theme
    _current_presentation = Presentation()

    title_slide_layout = _current_presentation.slide_layouts[0]
    slide = _current_presentation.slides.add_slide(title_slide_layout)

    # Set the title text
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Apply theme to title slide
    _theme_title_slide(slide, _active_theme)

    msg = f"Created presentation: '{title}' | theme='{_active_theme.name}'"
    logger.info(msg)
    return ToolResult(result=msg)


# ═══════════════════════════════════════════════════════════════════════════
# MCP TOOL: add_slide
# ═══════════════════════════════════════════════════════════════════════════

@mcp.tool()
def add_slide(title: str) -> SlideToolResult:
    """
    Adds a new themed content slide to the current presentation.

    Theme styling is applied automatically — do NOT try to style slides
    from the agent side; all styling lives here.
    """
    global _current_presentation, _active_theme

    if _current_presentation is None:
        return SlideToolResult(slide_index=-1,
                               result="Error: No presentation created yet.")

    # Pick best layout (prefer layout[1] = Title + Content)
    layout_to_use = None
    if len(_current_presentation.slide_layouts) > 1:
        layout_to_use = _current_presentation.slide_layouts[1]
    else:
        for layout in _current_presentation.slide_layouts:
            if len(layout.placeholders) >= 2:
                layout_to_use = layout
                break

    if layout_to_use is None:
        layout_to_use = _current_presentation.slide_layouts[0]

    slide = _current_presentation.slides.add_slide(layout_to_use)

    # Set title text
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Apply theme
    _theme_content_slide(slide, _active_theme)

    slide_index = len(_current_presentation.slides) - 1
    logger.info(f"[SLIDE] Added slide {slide_index}: '{title}' | theme='{_active_theme.name}'")
    return SlideToolResult(slide_index=slide_index,
                           result=f"Slide {slide_index} added with theme '{_active_theme.name}'")


# ═══════════════════════════════════════════════════════════════════════════
# MCP TOOL: write_text
# ═══════════════════════════════════════════════════════════════════════════

@mcp.tool()
def write_text(
    slide_index: Annotated[int, "The 0-indexed position of the slide"],
    points: Annotated[List[str], "A list of strings to be added as bullet points"],
) -> SlideToolResult:
    """
    Writes bullet points to a specific slide, applying theme body styling.
    """
    global _current_presentation, _active_theme

    # Robustness: unwrap dict payloads if the agent passes them
    if isinstance(points, dict) and "points" in points:
        points = points["points"]
    if isinstance(slide_index, dict) and "slide_index" in slide_index:
        slide_index = slide_index["slide_index"]

    if _current_presentation is None:
        return SlideToolResult(slide_index=-1,
                               result="Error: No presentation created yet.")

    if slide_index >= len(_current_presentation.slides):
        return SlideToolResult(slide_index=slide_index,
                               result=f"Error: Slide index {slide_index} out of range.")

    slide = _current_presentation.slides[slide_index]

    # Find the best body placeholder
    preferred_types = [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE]
    body_shape = None

    for p_type in preferred_types:
        for shape in slide.placeholders:
            if shape.placeholder_format.type == p_type:
                body_shape = shape
                break
        if body_shape:
            break

    if body_shape is None:
        for shape in slide.placeholders:
            if shape.placeholder_format.type != PP_PLACEHOLDER.TITLE:
                body_shape = shape
                break

    if body_shape is None:
        logger.warning(f"[WRITE] No placeholder on slide {slide_index}; creating textbox.")
        left, top, width, height = Inches(1), Inches(1.8), Inches(8), Inches(5)
        body_shape = slide.shapes.add_textbox(left, top, width, height)

    try:
        tf = body_shape.text_frame
        tf.word_wrap = True
        tf.text = ""

        for point in points:
            # Strip any [Tag] / (Tag) prefix — defence-in-depth cleaner
            clean_point = _LABEL_RE.sub("", str(point)).strip()
            if clean_point and clean_point[0].islower():
                clean_point = clean_point[0].upper() + clean_point[1:]
            clean_point = clean_point or str(point)   # fallback to original if empty

            p = tf.add_paragraph()
            p.text = clean_point
            p.level = 0
            # Apply theme body styling to each run
            for run in p.runs:
                run.font.color.rgb = _rgb(_active_theme.body_color)
                run.font.name      = _active_theme.font_family
                run.font.size      = Pt(_active_theme.body_size)
            # Also set paragraph-level defaults
            p.font.color.rgb = _rgb(_active_theme.body_color)
            p.font.name      = _active_theme.font_family
            p.font.size      = Pt(_active_theme.body_size)

        logger.info(f"[WRITE] {len(points)} bullet(s) → slide {slide_index}")
        return SlideToolResult(slide_index=slide_index,
                               result="Content added successfully")

    except Exception as exc:
        logger.error(f"[WRITE] Slide {slide_index} error: {exc}")
        return SlideToolResult(slide_index=slide_index, result=f"Error: {exc}")


# ═══════════════════════════════════════════════════════════════════════════
# MCP TOOL: save_presentation
# ═══════════════════════════════════════════════════════════════════════════

@mcp.tool()
def save_presentation(filename: str = "final.pptx") -> ToolResult:
    """
    Saves the current presentation to the output folder.
    """
    global _current_presentation

    if _current_presentation is None:
        return ToolResult(result="Error: No presentation to save.")

    os.makedirs(_output_path, exist_ok=True)
    full_path = os.path.join(_output_path, filename)
    _current_presentation.save(full_path)

    msg = f"Presentation saved to {full_path}"
    logger.info(msg)
    return ToolResult(result=msg)


# ── Entrypoint ─────────────────────────────────────────────────────────────
if __name__ == "__main__":
    mcp.run()
